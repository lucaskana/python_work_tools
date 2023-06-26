import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding an AutoShape'

left = Inches(0.93)  # 0.93" centers this overall set of shapes
#top = Inches(3.0)
top = Inches(0.93)
width = Inches(1.75)
height = Inches(1.0)

shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#shape.text = 'Step 1'

fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)
#fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
#fill.fore_color.brightness = 0.25
#fill.transparency = 0.25

#left = left + width - Inches(0.4)
top = top + height + Inches(0.2)
#width = Inches(2.0)  # chevrons need more width for visual balance

for n in range(2, 6):
    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.text = 'Step %d' % n
    top = top + height + Inches(0.2)

prs.save('data/arquivos/test.pptx')